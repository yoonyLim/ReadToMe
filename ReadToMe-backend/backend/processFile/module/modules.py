from io import BytesIO

from PyPDF2 import PdfReader
import olefile
import zlib
import struct
import docx2txt
from pptx import Presentation
from PIL import Image
import pytesseract

from gtts import gTTS

import re

# reference: https://blog.naver.com/PostView.naver?blogId=shino1025&logNo=222131312454&parentCategoryNo=&categoryNo=14&viewDate=&isShowPopularPosts=false&from=postView#post_1
def hwpToTxt(file):
    f = olefile.OleFileIO(file)
    dirs = f.listdir()
    text = ""

    # check if compressed
    header = f.openstream("FileHeader")
    header_data = header.read()
    is_compressed = (header_data[36] & 1) == 1

    nums = []
    for d in dirs:
        if d[0] == "BodyText":
            nums.append(int(d[1][len("Section"):]))
    sections = ["BodyText/Section"+str(x) for x in sorted(nums)]

    for section in sections:
        bodytext = f.openstream(section)
        data = bodytext.read()

        if is_compressed:
            data = zlib.decompress(data, -15)

        # extract txt from each section
        sectoin_text = ""
        i = 0
        size = len(data)

        while i < size:
            header = struct.unpack_from("<I", data, i)[0]
            rec_type = header & 0x3ff
            rec_len = (header >> 20) & 0xfff

            if rec_type in [67]:
                rec_data = data[i + 4 : i + 4 + rec_len]
                sectoin_text += rec_data.decode('utf-16')
                sectoin_text += '\n'

            i += 4 + rec_len

        text += sectoin_text
        text += '\n'

    return text

# reference: https://wooiljeong.github.io/python/pdf-to-text/
def pdfToTxt(file):
    reader = PdfReader(file)
    pages = reader.pages
    text = ""

    for page in pages:
        sub = page.extract_text()
        text += sub

    return text

def docxToTxt(file):
    return docx2txt.process(file)

def pptxToTxt(file):
    f = Presentation(file)
    text = ""

    for slide in f.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text

    return text

def imgToTxt(file):
    return pytesseract.image_to_string(Image.open(file), lang = 'kor+eng', timeout = 5)

def txtToSpch(txt):
    txt = re.sub(r"[^\uAC00-\uD7A30-9a-zA-Z\s@]", "", txt)

    tts_ko = gTTS(text = txt, lang = 'ko')
    tts_ko.save('result.mp3')